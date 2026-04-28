"""
Creates enhanced version of NID.pptx — keeps existing 7 slides,
inserts 4 new detailed slides after Slide 4 (Dataset):
  - Slide 5: Dataset Deep Dive (features breakdown, sample row)
  - Slide 6: Attack Types Explained (visual cards for each attack)
  - Slide 7: Data Preprocessing Pipeline (step-by-step visual flow)
  - Slide 8: Feature Engineering & SMOTE (visual diagrams)
Then original slides 5-7 become slides 9-11.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Theme colors extracted from user's PPT ──
TEAL = RGBColor(0x00, 0x89, 0x8A)
NAVY = RGBColor(0x1A, 0x2A, 0x4A)
CORAL = RGBColor(0xE8, 0x63, 0x52)
AMBER = RGBColor(0xE0, 0x9F, 0x1F)
SLATE = RGBColor(0x47, 0x55, 0x69)
SOFT_GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_WHITE = RGBColor(0xFA, 0xFA, 0xFC)
CARD_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEAL = RGBColor(0xF0, 0xFD, 0xFA)
LIGHT_CORAL = RGBColor(0xFE, 0xF2, 0xF2)
LIGHT_AMBER = RGBColor(0xFF, 0xFB, 0xEB)
LIGHT_NAVY = RGBColor(0xEF, 0xF6, 0xFF)
LIGHT_SLATE = RGBColor(0xF8, 0xFA, 0xFC)
BORDER_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
LIGHT_PURPLE = RGBColor(0xF5, 0xF3, 0xFF)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
LIGHT_GREEN = RGBColor(0xF0, 0xFD, 0xF4)

# Slide dimensions from user's PPT (standard 10x5.63)
SW = 9144000  # EMU
SH = 5143500

def emu(inches):
    return int(inches * 914400)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_left_bar(slide):
    """Teal left accent bar matching user's theme"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, emu(0.09), SH)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()

def add_header(slide, title_text):
    """Standard header with title + teal underline"""
    tb = slide.shapes.add_textbox(emu(0.6), emu(0.3), emu(8.3), emu(0.48))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.color.rgb = NAVY
    p.font.bold = True
    p.font.name = "Calibri"

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(0.6), emu(0.82), emu(2.25), emu(0.037))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

def add_card(slide, left, top, width, height, fill_color=CARD_WHITE, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.2)
    else:
        shape.line.color.rgb = BORDER_GRAY
        shape.line.width = Pt(0.75)
    return shape

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=14, color=SLATE, bold=False, align=PP_ALIGN.LEFT, name="Calibri"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return p

def add_para(tf, text, size=14, color=SLATE, bold=False, space_before=Pt(3), align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = space_before
    p.space_after = Pt(1)
    p.alignment = align
    return p

def add_rich_bullet(tf, label, desc, label_color=TEAL, desc_color=SLATE, size=11, space_before=Pt(4)):
    """Add a bullet with bold colored label + normal description"""
    p = tf.add_paragraph()
    p.space_before = space_before
    p.space_after = Pt(1)
    run1 = p.add_run()
    run1.text = label + ": "
    run1.font.size = Pt(size)
    run1.font.color.rgb = label_color
    run1.font.bold = True
    run1.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(size)
    run2.font.color.rgb = desc_color
    run2.font.bold = False
    run2.font.name = "Calibri"
    return p

def add_arrow_shape(slide, left, top, width, height, color=TEAL):
    """Right arrow"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_chevron(slide, left, top, width, height, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_step_box(slide, x, y, w, h, number, title, desc, accent_color, bg_color):
    """Numbered step card with accent top border"""
    # Card
    card = add_card(slide, x, y, w, h, bg_color, accent_color)

    # Number circle
    circle = add_circle(slide, x + emu(0.12), y + emu(0.1), emu(0.32), accent_color)
    num_tb = add_textbox(slide, x + emu(0.12), y + emu(0.1), emu(0.32), emu(0.32))
    set_text(num_tb.text_frame, str(number), size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Title
    ttb = add_textbox(slide, x + emu(0.5), y + emu(0.1), w - emu(0.6), emu(0.3))
    set_text(ttb.text_frame, title, size=12, color=NAVY, bold=True)

    # Description
    dtb = add_textbox(slide, x + emu(0.12), y + emu(0.45), w - emu(0.24), h - emu(0.55))
    return dtb


# ══════════════════════════════════════════════
# Build the presentation
# ══════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH


# ══════════════════════════════════════════════
# NEW SLIDE A: Dataset Deep Dive — Feature Categories
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
add_left_bar(slide)
add_header(slide, "Dataset Deep Dive — The 41 Features")

# 4 feature category cards
categories = [
    ("Basic Features", "1 — 9", "Direct TCP/IP attributes", [
        ("duration", "Connection length in seconds"),
        ("protocol_type", "TCP, UDP, or ICMP"),
        ("service", "HTTP, FTP, SMTP, etc. (70 types)"),
        ("flag", "Connection status (SF, REJ, S0...)"),
        ("src_bytes", "Bytes sent by source"),
        ("dst_bytes", "Bytes received by destination"),
    ], TEAL, LIGHT_TEAL),
    ("Content Features", "10 — 22", "Payload & access attributes", [
        ("num_failed_logins", "Failed login attempts"),
        ("logged_in", "Successfully logged in (1/0)"),
        ("root_shell", "Root shell obtained (1/0)"),
        ("num_file_creations", "File creation operations"),
        ("num_access_files", "Sensitive file accesses"),
        ("is_guest_login", "Guest login detected"),
    ], CORAL, LIGHT_CORAL),
    ("Time-based Traffic", "23 — 31", "2-second window statistics", [
        ("count", "Connections to same host"),
        ("srv_count", "Same-service connections"),
        ("serror_rate", "% connections with SYN errors"),
        ("rerror_rate", "% connections with REJ errors"),
        ("same_srv_rate", "% to same service"),
        ("diff_srv_rate", "% to different services"),
    ], AMBER, LIGHT_AMBER),
    ("Host-based Traffic", "32 — 41", "Last 100 connections stats", [
        ("dst_host_count", "Connections to same dest"),
        ("dst_host_srv_count", "Same-service to dest"),
        ("dst_host_same_srv_rate", "% same service to host"),
        ("dst_host_serror_rate", "% SYN errors to host"),
        ("dst_host_srv_serror_rate", "% SYN errors to service"),
        ("dst_host_rerror_rate", "% REJ errors to host"),
    ], NAVY, LIGHT_NAVY),
]

for i, (cat_title, feat_range, subtitle, features, accent, bg) in enumerate(categories):
    x = emu(0.45) + emu(2.35) * i
    y = emu(1.05)
    w = emu(2.2)
    h = emu(4.3)

    card = add_card(slide, x, y, w, h, bg, accent)

    # Category header with colored bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, emu(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # Title
    ttb = add_textbox(slide, x + emu(0.1), y + emu(0.1), w - emu(0.2), emu(0.35))
    set_text(ttb.text_frame, cat_title, size=12, color=accent, bold=True)
    add_para(ttb.text_frame, f"Features {feat_range}", size=9, color=SOFT_GRAY, space_before=Pt(1))

    # Subtitle
    stb = add_textbox(slide, x + emu(0.1), y + emu(0.5), w - emu(0.2), emu(0.22))
    set_text(stb.text_frame, subtitle, size=9, color=SLATE, bold=False)

    # Feature list
    ftb = add_textbox(slide, x + emu(0.1), y + emu(0.75), w - emu(0.2), h - emu(0.85))
    first = True
    for fname, fdesc in features:
        if first:
            set_text(ftb.text_frame, "", size=9, color=SLATE)
            first = False
        p = ftb.text_frame.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(0)
        r1 = p.add_run()
        r1.text = fname
        r1.font.size = Pt(9)
        r1.font.color.rgb = accent
        r1.font.bold = True
        r1.font.name = "Consolas"
        p2 = ftb.text_frame.add_paragraph()
        p2.space_before = Pt(0)
        p2.space_after = Pt(1)
        r2 = p2.add_run()
        r2.text = fdesc
        r2.font.size = Pt(8)
        r2.font.color.rgb = SOFT_GRAY
        r2.font.name = "Calibri"

# Bottom note
ntb = add_textbox(slide, emu(0.5), emu(5.45), emu(9), emu(0.2))
set_text(ntb.text_frame, "3 categorical features (protocol_type, service, flag) need encoding  |  38 numerical features need scaling  |  Last column = difficulty score (dropped)",
         size=8, color=SOFT_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════
# NEW SLIDE B: Attack Types Explained (Visual)
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
add_left_bar(slide)
add_header(slide, "Understanding the 5 Traffic Classes")

attacks = [
    ("Normal", "53.46%", "67,343", "Legitimate network traffic. Regular HTTP requests, file transfers, emails — nothing suspicious.",
     "shield", GREEN, LIGHT_GREEN),
    ("DoS", "36.46%", "45,927", "Denial of Service — flood the server with traffic until it crashes. Like 10,000 calls to a restaurant at once. Real users can't get through.",
     "flash_on", CORAL, LIGHT_CORAL),
    ("Probe", "9.25%", "11,656", "Reconnaissance scanning — checking which ports are open, what services run. Not attacking yet, but mapping vulnerabilities for later.",
     "search", AMBER, LIGHT_AMBER),
    ("R2L", "0.79%", "995", "Remote to Local — attacker is outside trying to get in. Password guessing, exploiting FTP, trying to gain unauthorized access.",
     "lock_open", PURPLE, LIGHT_PURPLE),
    ("U2R", "0.04%", "52", "User to Root — already inside with basic access, trying to become admin. Buffer overflow exploits, rootkits, privilege escalation.",
     "admin_panel_settings", NAVY, LIGHT_NAVY),
]

for i, (name, pct, count, desc, icon, accent, bg) in enumerate(attacks):
    y = emu(1.0) + emu(0.85) * i

    # Full-width card
    card = add_card(slide, emu(0.45), y, emu(9.1), emu(0.78), bg, accent)

    # Left accent bar on card
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(0.45), y, emu(0.06), emu(0.78))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # Class name
    ntb = add_textbox(slide, emu(0.65), y + emu(0.08), emu(0.8), emu(0.32))
    set_text(ntb.text_frame, name, size=16, color=accent, bold=True)

    # Percentage + count
    ptb = add_textbox(slide, emu(0.65), y + emu(0.4), emu(0.9), emu(0.3))
    set_text(ptb.text_frame, f"{pct}  ({count})", size=9, color=SOFT_GRAY, bold=False)

    # Description
    dtb = add_textbox(slide, emu(1.65), y + emu(0.05), emu(7.7), emu(0.68))
    set_text(dtb.text_frame, desc, size=10, color=SLATE)

    # Danger indicator (visual bar showing proportion)
    bar_width = float(pct.replace('%','')) / 53.46 * emu(0.5)  # scale relative to Normal
    pbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(9.2), y + emu(0.3), max(int(bar_width), emu(0.04)), emu(0.15))
    pbar.fill.solid()
    pbar.fill.fore_color.rgb = accent
    pbar.line.fill.background()

# Bottom insight box
icard = add_card(slide, emu(0.45), emu(5.35), emu(9.1), emu(0.4), LIGHT_CORAL, CORAL)
itb = add_textbox(slide, emu(0.65), emu(5.38), emu(8.7), emu(0.35))
tf = itb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r1 = p.add_run()
r1.text = "Key Insight: "
r1.font.size = Pt(9)
r1.font.color.rgb = CORAL
r1.font.bold = True
r1.font.name = "Calibri"
r2 = p.add_run()
r2.text = "R2L + U2R together = less than 1% of data, but these are the most dangerous attacks (unauthorized access + privilege escalation). Detecting them is the real challenge."
r2.font.size = Pt(9)
r2.font.color.rgb = SLATE
r2.font.bold = False
r2.font.name = "Calibri"


# ══════════════════════════════════════════════
# NEW SLIDE C: Data Preprocessing Pipeline
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
add_left_bar(slide)
add_header(slide, "Data Preprocessing Pipeline")

# Step 1: Load & Clean
dtb = add_step_box(slide, emu(0.4), emu(1.0), emu(2.9), emu(1.75), 1, "Load & Inspect", "", TEAL, LIGHT_TEAL)
set_text(dtb.text_frame, "", size=9, color=SLATE)
add_rich_bullet(dtb.text_frame, "Load", "KDDTrain+.txt (125,973 rows)", TEAL, SLATE, size=9)
add_rich_bullet(dtb.text_frame, "Assign", "41 column names + label", TEAL, SLATE, size=9)
add_rich_bullet(dtb.text_frame, "Drop", "difficulty score (last col)", TEAL, SLATE, size=9)
add_rich_bullet(dtb.text_frame, "Map", "specific attacks to 5 classes (e.g., neptune, smurf -> DoS)", TEAL, SLATE, size=9)

# Arrow 1
add_arrow_shape(slide, emu(3.42), emu(1.7), emu(0.32), emu(0.25), TEAL)

# Step 2: Encode
dtb = add_step_box(slide, emu(3.85), emu(1.0), emu(2.9), emu(1.75), 2, "Encode Categoricals", "", CORAL, LIGHT_CORAL)
set_text(dtb.text_frame, "", size=9, color=SLATE)
add_rich_bullet(dtb.text_frame, "protocol_type", "3 values (TCP, UDP, ICMP)", CORAL, SLATE, size=9)
add_rich_bullet(dtb.text_frame, "service", "70 values (HTTP, FTP...)", CORAL, SLATE, size=9)
add_rich_bullet(dtb.text_frame, "flag", "11 values (SF, REJ, S0...)", CORAL, SLATE, size=9)
add_rich_bullet(dtb.text_frame, "Method", "One-Hot Encoding (no false ordering)", CORAL, SLATE, size=9)

# Arrow 2
add_arrow_shape(slide, emu(6.87), emu(1.7), emu(0.32), emu(0.25), CORAL)

# Step 3: Scale
dtb = add_step_box(slide, emu(7.3), emu(1.0), emu(2.6), emu(1.75), 3, "Normalize / Scale", "", AMBER, LIGHT_AMBER)
set_text(dtb.text_frame, "", size=9, color=SLATE)
add_rich_bullet(dtb.text_frame, "Why", "Features have different scales (duration: 0-58000 vs failed_logins: 0-5)", AMBER, SLATE, size=9)
add_rich_bullet(dtb.text_frame, "Method", "StandardScaler (mean=0, std=1)", AMBER, SLATE, size=9)
add_rich_bullet(dtb.text_frame, "Critical for", "KNN & SVM (distance-based)", AMBER, SLATE, size=9)

# Row 2: Steps 4-5

# Step 4: SMOTE
dtb = add_step_box(slide, emu(0.4), emu(3.1), emu(4.6), emu(2.3), 4, "Handle Class Imbalance (SMOTE)", "", PURPLE, LIGHT_PURPLE)
set_text(dtb.text_frame, "", size=9, color=SLATE)

# Visual: Before vs After
p = dtb.text_frame.add_paragraph()
p.space_before = Pt(2)
r = p.add_run()
r.text = "Before SMOTE:"
r.font.size = Pt(9)
r.font.color.rgb = NAVY
r.font.bold = True
r.font.name = "Calibri"

for cls, cnt, bar_len in [("Normal", "67,343", 30), ("DoS", "45,927", 20), ("Probe", "11,656", 5), ("R2L", "995", 1), ("U2R", "52", 0)]:
    p2 = dtb.text_frame.add_paragraph()
    p2.space_before = Pt(1)
    r2 = p2.add_run()
    r2.text = f"  {cls:8s} {cnt:>8s}  {'█' * max(bar_len,1)}"
    r2.font.size = Pt(8)
    r2.font.color.rgb = PURPLE if bar_len <= 1 else SLATE
    r2.font.bold = bar_len <= 1
    r2.font.name = "Consolas"

p3 = dtb.text_frame.add_paragraph()
p3.space_before = Pt(5)
r3 = p3.add_run()
r3.text = "SMOTE creates synthetic minority samples by interpolating"
r3.font.size = Pt(9)
r3.font.color.rgb = PURPLE
r3.font.bold = False
r3.font.name = "Calibri"
p4 = dtb.text_frame.add_paragraph()
p4.space_before = Pt(1)
r4 = p4.add_run()
r4.text = "between existing samples and their nearest neighbors."
r4.font.size = Pt(9)
r4.font.color.rgb = PURPLE
r4.font.bold = False
r4.font.name = "Calibri"
p5 = dtb.text_frame.add_paragraph()
p5.space_before = Pt(3)
r5 = p5.add_run()
r5.text = "Applied ONLY to training folds, never to test data."
r5.font.size = Pt(9)
r5.font.color.rgb = CORAL
r5.font.bold = True
r5.font.name = "Calibri"

# Arrow
add_arrow_shape(slide, emu(5.1), emu(4.05), emu(0.32), emu(0.25), PURPLE)

# Step 5: Split
dtb = add_step_box(slide, emu(5.55), emu(3.1), emu(4.35), emu(2.3), 5, "Train-Test Split & Cross-Validation", "", NAVY, LIGHT_NAVY)
set_text(dtb.text_frame, "", size=9, color=SLATE)
add_rich_bullet(dtb.text_frame, "Primary split", "KDDTrain+ (125K) vs KDDTest+ (22K) — predefined", NAVY, SLATE, size=9)
add_rich_bullet(dtb.text_frame, "Cross-validation", "Stratified 10-fold on training set", NAVY, SLATE, size=9)
add_rich_bullet(dtb.text_frame, "Stratified", "Each fold preserves original class ratios", NAVY, SLATE, size=9)

# Visual: 10-fold illustration
p = dtb.text_frame.add_paragraph()
p.space_before = Pt(6)
r = p.add_run()
r.text = "10-Fold Cross-Validation:"
r.font.size = Pt(9)
r.font.color.rgb = NAVY
r.font.bold = True
r.font.name = "Calibri"

for fold in range(1, 4):
    p2 = dtb.text_frame.add_paragraph()
    p2.space_before = Pt(1)
    blocks = ""
    for b in range(1, 11):
        if b == fold:
            blocks += " [TEST]"
        else:
            blocks += " [TRAIN]"
    r2 = p2.add_run()
    r2.text = f"  Fold {fold}:{blocks}"
    r2.font.size = Pt(7)
    r2.font.color.rgb = SLATE
    r2.font.name = "Consolas"

p3 = dtb.text_frame.add_paragraph()
p3.space_before = Pt(1)
r3 = p3.add_run()
r3.text = "  ...    (repeated 10 times, each fold tested once)"
r3.font.size = Pt(7)
r3.font.color.rgb = SOFT_GRAY
r3.font.name = "Consolas"


# ══════════════════════════════════════════════
# NEW SLIDE D: Feature Engineering & Model Evaluation Strategy
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
add_left_bar(slide)
add_header(slide, "Feature Engineering & Evaluation Strategy")

# Left: Feature Engineering
card = add_card(slide, emu(0.4), emu(1.0), emu(4.65), emu(4.6), CARD_WHITE, TEAL)

# FE Header
ftb = add_textbox(slide, emu(0.55), emu(1.08), emu(4.3), emu(0.3))
set_text(ftb.text_frame, "Feature Engineering Steps", size=14, color=TEAL, bold=True)

# FE Content
ftb2 = add_textbox(slide, emu(0.55), emu(1.4), emu(4.3), emu(4.0))
set_text(ftb2.text_frame, "", size=9, color=SLATE)

# Step A
add_rich_bullet(ftb2.text_frame, "A. Correlation Analysis", "", TEAL, SLATE, size=10, space_before=Pt(4))
add_para(ftb2.text_frame, "Identify highly correlated feature pairs (r > 0.9). Remove one from each pair to reduce redundancy without losing information.", size=9, color=SLATE, space_before=Pt(2))

# Step B
add_rich_bullet(ftb2.text_frame, "B. Feature Importance Ranking", "", TEAL, SLATE, size=10, space_before=Pt(8))
add_para(ftb2.text_frame, "Use Random Forest / XGBoost built-in importance scores to rank features by their contribution to classification.", size=9, color=SLATE, space_before=Pt(2))

# Step C
add_rich_bullet(ftb2.text_frame, "C. Variance Threshold", "", TEAL, SLATE, size=10, space_before=Pt(8))
add_para(ftb2.text_frame, "Remove near-zero variance features that carry no useful information for distinguishing classes.", size=9, color=SLATE, space_before=Pt(2))

# Expected insight box
ibox = add_card(slide, emu(0.55), emu(4.15), emu(4.35), emu(1.2), LIGHT_TEAL, TEAL)
itb = add_textbox(slide, emu(0.7), emu(4.2), emu(4.0), emu(1.1))
tf = itb.text_frame
tf.word_wrap = True
set_text(tf, "Expected Finding", size=10, color=TEAL, bold=True)
add_para(tf, "Literature suggests ~10-15 features out of 41 carry most discriminative power. Key features likely include:", size=9, color=SLATE, space_before=Pt(3))
add_para(tf, "src_bytes, dst_bytes, service, flag, count, srv_count, same_srv_rate, dst_host_srv_count", size=8, color=TEAL, bold=True, space_before=Pt(2))

# Right: Evaluation Metrics
card2 = add_card(slide, emu(5.25), emu(1.0), emu(4.65), emu(4.6), CARD_WHITE, CORAL)

etb = add_textbox(slide, emu(5.4), emu(1.08), emu(4.3), emu(0.3))
set_text(etb.text_frame, "Evaluation Metrics Explained", size=14, color=CORAL, bold=True)

metrics = [
    ("Accuracy", "Overall correct predictions / total. Misleading with imbalanced data — a model predicting 'Normal' for everything gets 53%.", CORAL),
    ("Precision", "Of all predicted attacks, how many were real? High precision = few false alarms.", AMBER),
    ("Recall", "Of all actual attacks, how many were caught? High recall = few missed attacks. Critical for security.", GREEN),
    ("F1-Score", "Harmonic mean of Precision & Recall. Both must be good for F1 to be high. Our primary metric.", NAVY),
    ("ROC-AUC", "Trade-off between true positive & false positive rates. Closer to 1.0 = better. 0.5 = random guess.", PURPLE),
    ("Confusion Matrix", "5x5 grid showing exactly which classes get confused with which. Reveals where each model fails.", TEAL),
]

mtb = add_textbox(slide, emu(5.4), emu(1.4), emu(4.3), emu(4.1))
set_text(mtb.text_frame, "", size=9, color=SLATE)

for mname, mdesc, mcolor in metrics:
    add_rich_bullet(mtb.text_frame, mname, mdesc, mcolor, SLATE, size=9, space_before=Pt(5))


# ══════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════
output = "/Users/sudipadh/Desktop/masters/sem 1/machine_learning/docs/NID_enhanced_slides.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
