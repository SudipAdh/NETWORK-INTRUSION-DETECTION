from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Light theme color scheme - clean, professional
BG_WHITE = RGBColor(0xFA, 0xFA, 0xFC)
NAVY = RGBColor(0x1A, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0x89, 0x8A)
CORAL = RGBColor(0xE8, 0x63, 0x52)
AMBER = RGBColor(0xE0, 0x9F, 0x1F)
SLATE = RGBColor(0x47, 0x55, 0x69)
SOFT_GRAY = RGBColor(0x64, 0x74, 0x8B)
BORDER_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
CARD_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEAL_BG = RGBColor(0xF0, 0xFD, 0xFA)
LIGHT_CORAL_BG = RGBColor(0xFE, 0xF2, 0xF2)
LIGHT_AMBER_BG = RGBColor(0xFF, 0xFB, 0xEB)
LIGHT_NAVY_BG = RGBColor(0xEF, 0xF6, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(text_frame, text, font_size=18, color=NAVY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    text_frame.clear()
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return p

def add_paragraph(text_frame, text, font_size=18, color=NAVY, bold=False, space_before=Pt(6), space_after=Pt(2), font_name="Calibri", alignment=PP_ALIGN.LEFT):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    return p

def add_bullet(text_frame, text, font_size=16, color=SLATE, level=0, bold=False, space_before=Pt(4)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.level = level
    p.space_before = space_before
    p.space_after = Pt(2)
    return p


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)

# Top accent bar
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), TEAL)

# Left accent stripe
add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), TEAL)

# Title
tb = add_text_box(slide, Inches(1), Inches(1.3), Inches(11.3), Inches(2))
set_text(tb.text_frame, "A Comparative Study of Classical Machine Learning", font_size=38, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(tb.text_frame, "Algorithms for Multi-Class Network Intrusion Detection", font_size=38, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
tb2 = add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.6))
set_text(tb2.text_frame, "Using the NSL-KDD Benchmark Dataset", font_size=22, color=TEAL, bold=False, alignment=PP_ALIGN.CENTER)

# Divider
add_rect(slide, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.04), TEAL)

# Author info
tb3 = add_text_box(slide, Inches(1), Inches(4.7), Inches(11.3), Inches(2.2))
set_text(tb3.text_frame, "Sudip Adhikari (250578)", font_size=24, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(tb3.text_frame, "M.Sc. Computer Engineering", font_size=17, color=SLATE, alignment=PP_ALIGN.CENTER, space_before=Pt(8))
add_paragraph(tb3.text_frame, "Softwarica College of IT & E-Commerce", font_size=17, color=SLATE, alignment=PP_ALIGN.CENTER, space_before=Pt(4))
add_paragraph(tb3.text_frame, "in collaboration with Coventry University", font_size=15, color=SOFT_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(4))
add_paragraph(tb3.text_frame, "Machine Learning  |  Project Proposal  |  2026", font_size=14, color=SOFT_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(20))


# ============================================================
# SLIDE 2: Why This Topic?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), TEAL)

# Header
tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
set_text(tb.text_frame, "Why Network Intrusion Detection?", font_size=32, color=NAVY, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.05), TEAL)

# Left column - The Problem
card1 = add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2), LIGHT_NAVY_BG, TEAL)
tb1 = add_text_box(slide, Inches(1.2), Inches(1.8), Inches(4.9), Inches(4.8))
set_text(tb1.text_frame, "The Problem", font_size=22, color=TEAL, bold=True)
add_bullet(tb1.text_frame, "Cyberattacks are growing fast globally and in Nepal — most organizations still rely on basic signature-based firewalls", font_size=15, color=SLATE, space_before=Pt(14))
add_bullet(tb1.text_frame, "Traditional Intrusion Detection Systems (IDS) use predefined rules — they catch known attacks but completely miss new or modified ones", font_size=15, color=SLATE, space_before=Pt(10))
add_bullet(tb1.text_frame, "We need systems that can learn patterns from network traffic and flag suspicious activity automatically", font_size=15, color=SLATE, space_before=Pt(10))
add_bullet(tb1.text_frame, "Most existing studies test only 1-2 algorithms — there is no clear picture of which classical ML approach works best across different attack types", font_size=15, color=SLATE, space_before=Pt(10))

# Right column - Why I Picked This
card2 = add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.2), LIGHT_TEAL_BG, TEAL)
tb2 = add_text_box(slide, Inches(7.2), Inches(1.8), Inches(4.9), Inches(4.8))
set_text(tb2.text_frame, "Why I Picked This", font_size=22, color=NAVY, bold=True)
add_bullet(tb2.text_frame, "The NSL-KDD dataset is a gold-standard benchmark — publicly available, well-documented, no data collection headaches", font_size=15, color=SLATE, space_before=Pt(14))
add_bullet(tb2.text_frame, "Cybersecurity + ML is a growing intersection with real-world demand, both in industry and research", font_size=15, color=SLATE, space_before=Pt(10))
add_bullet(tb2.text_frame, "The project is scoped right — not a weekend homework, but not an impossible research task either", font_size=15, color=SLATE, space_before=Pt(10))
add_bullet(tb2.text_frame, "Nobody else in the class picked a cybersecurity topic — this gives a different perspective from the agriculture and IoT projects most people are doing", font_size=15, color=SLATE, space_before=Pt(10))


# ============================================================
# SLIDE 3: Objectives
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), TEAL)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
set_text(tb.text_frame, "Project Objectives", font_size=32, color=NAVY, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.05), TEAL)

objectives = [
    ("01", "Compare 6 Classical ML Algorithms", "Train and evaluate Logistic Regression, KNN, Decision Tree, Random Forest, SVM, and XGBoost on the same dataset under identical conditions", TEAL, LIGHT_TEAL_BG),
    ("02", "Multi-Class Attack Classification", "Classify network traffic into 5 categories: Normal, DoS, Probe, R2L, and U2R — not just binary normal-vs-attack", CORAL, LIGHT_CORAL_BG),
    ("03", "Handle Class Imbalance", "R2L and U2R attacks make up less than 2% of the data — explore techniques like SMOTE and class weighting to deal with this realistically", AMBER, LIGHT_AMBER_BG),
    ("04", "Proper Evaluation Beyond Accuracy", "Use precision, recall, F1-score, confusion matrices, and ROC-AUC — because accuracy alone is misleading when classes are imbalanced", NAVY, LIGHT_NAVY_BG),
]

for i, (num, title, desc, color, bg) in enumerate(objectives):
    x = Inches(0.8) if i % 2 == 0 else Inches(6.8)
    y = Inches(1.5) + Inches(2.7) * (i // 2)
    card = add_shape(slide, x, y, Inches(5.6), Inches(2.3), bg, color)

    num_tb = add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.5))
    set_text(num_tb.text_frame, num, font_size=22, color=color, bold=True)

    title_tb = add_text_box(slide, x + Inches(0.9), y + Inches(0.25), Inches(4.3), Inches(0.5))
    set_text(title_tb.text_frame, title, font_size=18, color=NAVY, bold=True)

    desc_tb = add_text_box(slide, x + Inches(0.4), y + Inches(0.9), Inches(4.8), Inches(1.2))
    set_text(desc_tb.text_frame, desc, font_size=14, color=SLATE)


# ============================================================
# SLIDE 4: Dataset
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), TEAL)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
set_text(tb.text_frame, "The Dataset — NSL-KDD", font_size=32, color=NAVY, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.05), TEAL)

# Left - About
card1 = add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(3), LIGHT_NAVY_BG, BORDER_GRAY)
tb1 = add_text_box(slide, Inches(1.2), Inches(1.7), Inches(4.9), Inches(2.6))
set_text(tb1.text_frame, "What is NSL-KDD?", font_size=20, color=TEAL, bold=True)
add_bullet(tb1.text_frame, "Most widely used benchmark for network intrusion detection research", font_size=15, color=SLATE, space_before=Pt(12))
add_bullet(tb1.text_frame, "Improved version of KDD Cup 1999 — removes duplicate records that used to bias results in older studies", font_size=15, color=SLATE, space_before=Pt(8))
add_bullet(tb1.text_frame, "Publicly available from the Canadian Institute for Cybersecurity — no access barriers, no data collection needed", font_size=15, color=SLATE, space_before=Pt(8))
add_bullet(tb1.text_frame, "41 features covering protocol type, service, flag, duration, bytes transferred, error rates, and more", font_size=15, color=SLATE, space_before=Pt(8))

# Right - Stats cards
stats = [
    ("Training Samples", "125,973", TEAL, LIGHT_TEAL_BG),
    ("Test Samples", "22,544", CORAL, LIGHT_CORAL_BG),
    ("Features", "41", AMBER, LIGHT_AMBER_BG),
    ("Classes", "5", NAVY, LIGHT_NAVY_BG),
]

for i, (label, value, color, bg) in enumerate(stats):
    x = Inches(6.8) + Inches(2.8) * (i % 2)
    y = Inches(1.5) + Inches(1.55) * (i // 2)
    card = add_shape(slide, x, y, Inches(2.5), Inches(1.3), bg, color)
    vtb = add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.1), Inches(0.6))
    set_text(vtb.text_frame, value, font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    ltb = add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.1), Inches(0.4))
    set_text(ltb.text_frame, label, font_size=13, color=SLATE, alignment=PP_ALIGN.CENTER)

# Bottom - Class distribution
card_bottom = add_shape(slide, Inches(0.8), Inches(4.8), Inches(11.6), Inches(2.1), LIGHT_CORAL_BG, CORAL)
tb_b = add_text_box(slide, Inches(1.2), Inches(4.95), Inches(10.8), Inches(1.8))
set_text(tb_b.text_frame, "Class Distribution (the challenge)", font_size=18, color=CORAL, bold=True)
add_bullet(tb_b.text_frame, "Normal: ~53%    |    DoS: ~36%    |    Probe: ~9%    |    R2L: ~1.7%    |    U2R: ~0.04%", font_size=15, color=NAVY, bold=True, space_before=Pt(12))
add_bullet(tb_b.text_frame, "R2L and U2R are extremely rare — this is realistic (rare attacks are rare in real life too) but makes classification much harder. Any model that just predicts 'Normal' for everything would still get ~53% accuracy. That is why we need proper metrics beyond just accuracy.", font_size=14, color=SLATE, space_before=Pt(8))


# ============================================================
# SLIDE 5: Methodology
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), TEAL)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
set_text(tb.text_frame, "Methodology", font_size=32, color=NAVY, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.05), TEAL)

# Pipeline steps
steps = [
    ("1. Data\nPreprocessing", "Encode categorical\nfeatures, normalize\nnumerical values,\nhandle imbalance\n(SMOTE)", TEAL, LIGHT_TEAL_BG),
    ("2. Feature\nEngineering", "Analyze feature\ncorrelations, remove\nredundant features,\nselect most\ninformative ones", CORAL, LIGHT_CORAL_BG),
    ("3. Model\nTraining", "Train 6 algorithms:\nLR, KNN, DT, RF,\nSVM, XGBoost\nusing 10-fold\ncross-validation", AMBER, LIGHT_AMBER_BG),
    ("4. Evaluation\n& Comparison", "Confusion matrices,\nper-class P/R/F1,\nROC-AUC curves,\nfeature importance\nranking", NAVY, LIGHT_NAVY_BG),
]

for i, (title, desc, color, bg) in enumerate(steps):
    x = Inches(0.6) + Inches(3.15) * i
    card = add_shape(slide, x, Inches(1.5), Inches(2.8), Inches(3.2), bg, color)
    ttb = add_text_box(slide, x + Inches(0.2), Inches(1.7), Inches(2.4), Inches(0.9))
    set_text(ttb.text_frame, title, font_size=17, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    dtb = add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.4), Inches(1.8))
    set_text(dtb.text_frame, desc, font_size=13, color=SLATE, alignment=PP_ALIGN.CENTER)

# Arrow connectors
for i in range(3):
    x = Inches(3.4) + Inches(3.15) * i
    arrow_tb = add_text_box(slide, x, Inches(2.7), Inches(0.5), Inches(0.5))
    set_text(arrow_tb.text_frame, ">", font_size=28, color=SOFT_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# Algorithms box at bottom
card_algo = add_shape(slide, Inches(0.8), Inches(5.1), Inches(11.6), Inches(1.8), LIGHT_TEAL_BG, TEAL)
tb_algo = add_text_box(slide, Inches(1.2), Inches(5.25), Inches(10.8), Inches(1.5))
set_text(tb_algo.text_frame, "Algorithms Being Compared", font_size=18, color=TEAL, bold=True)

algos = [
    ("Logistic Regression", "Linear baseline"),
    ("K-Nearest Neighbors", "Distance-based"),
    ("Decision Tree", "Rule-based, interpretable"),
    ("Random Forest", "Ensemble of trees"),
    ("SVM", "Margin-based classifier"),
    ("XGBoost", "Gradient boosting"),
]
algo_text = "    |    ".join([f"{name} ({desc})" for name, desc in algos])
add_paragraph(tb_algo.text_frame, algo_text, font_size=13, color=SLATE, space_before=Pt(12))


# ============================================================
# SLIDE 6: Expected Outcomes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), TEAL)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
set_text(tb.text_frame, "Expected Outcomes", font_size=32, color=NAVY, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.05), TEAL)

# What I expect
card1 = add_shape(slide, Inches(0.8), Inches(1.5), Inches(7.2), Inches(4), LIGHT_NAVY_BG, BORDER_GRAY)
tb1 = add_text_box(slide, Inches(1.2), Inches(1.7), Inches(6.5), Inches(3.6))
set_text(tb1.text_frame, "What I Expect to Find", font_size=20, color=NAVY, bold=True)
add_bullet(tb1.text_frame, "Tree-based models (Random Forest, XGBoost) will likely perform best overall — this is consistent with existing literature on NSL-KDD", font_size=15, color=SLATE, space_before=Pt(14))
add_bullet(tb1.text_frame, "The real interesting part is how each model handles rare classes (R2L, U2R) — some models might have good overall accuracy but completely fail on these", font_size=15, color=SLATE, space_before=Pt(8))
add_bullet(tb1.text_frame, "Simpler models like Logistic Regression may do well on binary detection (normal vs attack) but struggle with 5-class classification", font_size=15, color=SLATE, space_before=Pt(8))
add_bullet(tb1.text_frame, "Feature importance analysis will reveal which network traffic features are most useful for detecting intrusions — this has practical value for network security teams", font_size=15, color=SLATE, space_before=Pt(8))

# Deliverables
card2 = add_shape(slide, Inches(8.3), Inches(1.5), Inches(4.2), Inches(4), LIGHT_TEAL_BG, TEAL)
tb2 = add_text_box(slide, Inches(8.6), Inches(1.7), Inches(3.7), Inches(3.6))
set_text(tb2.text_frame, "Deliverables", font_size=20, color=TEAL, bold=True)
add_bullet(tb2.text_frame, "Complete Python codebase (reproducible)", font_size=14, color=SLATE, space_before=Pt(14))
add_bullet(tb2.text_frame, "Model comparison tables and charts", font_size=14, color=SLATE, space_before=Pt(8))
add_bullet(tb2.text_frame, "Per-class evaluation with confusion matrices", font_size=14, color=SLATE, space_before=Pt(8))
add_bullet(tb2.text_frame, "Feature importance rankings", font_size=14, color=SLATE, space_before=Pt(8))
add_bullet(tb2.text_frame, "Final report with analysis and discussion", font_size=14, color=SLATE, space_before=Pt(8))

# Bottom note
note = add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.6), Inches(1))
set_text(note.text_frame, "Note: These are hypotheses based on existing literature — the whole point of this project is to verify these with proper experiments and present honest findings, whether they match expectations or not.", font_size=13, color=SOFT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7: Timeline & Tools
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), TEAL)

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
set_text(tb.text_frame, "Timeline & Tools", font_size=32, color=NAVY, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.05), TEAL)

# 2-week timeline - split into day ranges
days = [
    ("Day 1-2", "Data loading, EDA,\npreprocessing,\nfeature engineering", TEAL, LIGHT_TEAL_BG),
    ("Day 3-5", "Train all 6 models:\nLR, KNN, DT, RF,\nSVM, XGBoost", CORAL, LIGHT_CORAL_BG),
    ("Day 6-8", "Evaluation, tuning,\nconfusion matrices,\nROC curves", AMBER, LIGHT_AMBER_BG),
    ("Day 9-11", "Feature importance,\ncomparison analysis,\nvisualizations", NAVY, LIGHT_NAVY_BG),
    ("Day 12-14", "Report writing,\nfinal presentation,\ncleanup", TEAL, LIGHT_TEAL_BG),
]

for i, (day, task, color, bg) in enumerate(days):
    x = Inches(0.5) + Inches(2.5) * i
    card = add_shape(slide, x, Inches(1.5), Inches(2.2), Inches(2.4), bg, color)
    wtb = add_text_box(slide, x + Inches(0.1), Inches(1.65), Inches(2.0), Inches(0.5))
    set_text(wtb.text_frame, day, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    ttb = add_text_box(slide, x + Inches(0.1), Inches(2.25), Inches(2.0), Inches(1.4))
    set_text(ttb.text_frame, task, font_size=14, color=SLATE, alignment=PP_ALIGN.CENTER)

# Week labels
w1_tb = add_text_box(slide, Inches(0.5), Inches(4.0), Inches(7.5), Inches(0.4))
set_text(w1_tb.text_frame, "Week 1", font_size=14, color=TEAL, bold=True, alignment=PP_ALIGN.LEFT)
add_rect(slide, Inches(1.5), Inches(4.15), Inches(4.0), Inches(0.03), BORDER_GRAY)

w2_tb = add_text_box(slide, Inches(5.5), Inches(4.0), Inches(7.5), Inches(0.4))
set_text(w2_tb.text_frame, "Week 2", font_size=14, color=CORAL, bold=True, alignment=PP_ALIGN.LEFT)
add_rect(slide, Inches(6.5), Inches(4.15), Inches(6.0), Inches(0.03), BORDER_GRAY)

# Tools section
card_tools = add_shape(slide, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.6), LIGHT_TEAL_BG, BORDER_GRAY)
tb_tools = add_text_box(slide, Inches(1.2), Inches(4.65), Inches(4.9), Inches(2.3))
set_text(tb_tools.text_frame, "Tools & Libraries", font_size=20, color=TEAL, bold=True)
add_bullet(tb_tools.text_frame, "Python 3.x — primary language", font_size=14, color=SLATE, space_before=Pt(10))
add_bullet(tb_tools.text_frame, "scikit-learn — ML algorithms & evaluation", font_size=14, color=SLATE, space_before=Pt(5))
add_bullet(tb_tools.text_frame, "XGBoost — gradient boosting implementation", font_size=14, color=SLATE, space_before=Pt(5))
add_bullet(tb_tools.text_frame, "pandas, NumPy — data handling", font_size=14, color=SLATE, space_before=Pt(5))
add_bullet(tb_tools.text_frame, "matplotlib, seaborn — visualizations", font_size=14, color=SLATE, space_before=Pt(5))
add_bullet(tb_tools.text_frame, "Jupyter Notebook — experimentation", font_size=14, color=SLATE, space_before=Pt(5))

# References section
card_ref = add_shape(slide, Inches(6.8), Inches(4.5), Inches(5.6), Inches(2.6), LIGHT_NAVY_BG, BORDER_GRAY)
tb_ref = add_text_box(slide, Inches(7.2), Inches(4.65), Inches(4.9), Inches(2.3))
set_text(tb_ref.text_frame, "Key References", font_size=20, color=NAVY, bold=True)
add_bullet(tb_ref.text_frame, "Tavallaee et al. (2009) — A detailed analysis of the KDD CUP 99 dataset", font_size=13, color=SLATE, space_before=Pt(10))
add_bullet(tb_ref.text_frame, "Ahmad et al. (2021) — Network intrusion detection using ML: A comparative study", font_size=13, color=SLATE, space_before=Pt(5))
add_bullet(tb_ref.text_frame, "Dhanabal & Shantharajah (2015) — A study on NSL-KDD dataset for IDS", font_size=13, color=SLATE, space_before=Pt(5))
add_bullet(tb_ref.text_frame, "Ingre & Yadav (2015) — Performance analysis of NSL-KDD using ANN", font_size=13, color=SLATE, space_before=Pt(5))

# Thank you footer
footer = add_text_box(slide, Inches(0.8), Inches(7.1), Inches(11.6), Inches(0.35))
set_text(footer.text_frame, "Thank You  —  Questions?", font_size=16, color=SOFT_GRAY, alignment=PP_ALIGN.CENTER)


# Save
output_path = "/Users/sudipadh/Desktop/masters/sem 1/machine_learning/docs/ML_Project_Proposal_Sudip_Adhikari.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
